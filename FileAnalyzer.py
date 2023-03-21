import os 
import re
import tkinter as tk
import PySimpleGUI as sg
from tkinter import filedialog
import nltk
import nltk.corpus
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
from sklearn.pipeline import Pipeline
from pptx import Presentation
import PyPDF2
import string
from nltk.tokenize import RegexpTokenizer
from nltk.stem import WordNetLemmatizer
from gensim import corpora, models
import gensim
import PySimpleGUI as sg
import pprint
nltk.download('stopwords')
nltk.download('punkt')
nltk.download("wordnet")
nltk.download("omw-1.4")


keyword1 = ['element', 'matrix', 'equation', 'function', 'boundary', 'solution', 'force', 'condition', 'form', 'finite', 'shape', 
            'elastic', 'model', 'node', 'method', 'lesson', 'stiffness', 'stress', 'number', 'displacement'] #Finite Element
keyword2 = ['employee', 'resource', 'human', 'competitive', 'management', 'manage', 'state', 'business', 'advantage'
            , 'gowan', 'introduction', 'performance', 'plan', 'organization', 'work'] #Human Resource
keyword3 = ['price', 'right', 'permit', 'service', 'distribute', 'license', 'certain', 'scan', 'reserve', 'copy', 
            'duplicate', 'produ', 'cengage'] #Economic
keywords = [keyword1, keyword2, keyword3]

class FileAnalyzer:
    def __init__(self, file_path, **kwargs):
        self.file_path = file_path
        self.wnl = WordNetLemmatizer()
        self.stop_words = set(stopwords.words("english"))
        self.punc =  ['!', '(', ')', '-', '[', ']', '{', '}', ';', ':', '\\', '<', '>', ',', '.', '/', '?', '@', '#', '$', '%',
                       '^', '&', '*', '_', '…', '~', '"""', '0', '1', '2', '3', '4', '5', '6', '7', '8', '9', '+', '=','||']
    
    def extract_text_from_ppt(self,file_path):
        prs = Presentation(file_path)
        text_list = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_list.append(shape.text)
        #convert all the letter to lower case
        for i in range(len(text_list)):
            text_list[i] = text_list[i].lower()    
        return text_list
    
    def extract_text_from_pdf(self, file_path):
        pdf_text_str = ''
        pdf_text_list = []
        pdf_reader = PyPDF2.PdfReader(file_path)
        num_pages = len(pdf_reader.pages)
        for page_num in range(num_pages):
            page_obj = pdf_reader.pages[page_num]
            pdf_text_str += page_obj.extract_text()
        #append the extract text to list
        pdf_text_list.append(pdf_text_str)
        #Convert all the letter to lower case
        for i in range(len(pdf_text_list)):
            pdf_text_list[i] = pdf_text_list[i].lower()
        return pdf_text_list
    
    def remove_stopwords(self,text):
        text_str= ''
        for t in text:
            text_str += ''+t
        #Tokenization
        words = nltk.tokenize.word_tokenize(text_str)
        #Remove stopwords
        filter_text1 = [w for w in words if w not in self.stop_words]
        #Remove punctation
        filter_text2 = []
        for ele in filter_text1:
            for a in self.punc:
                if a in ele:
                    x='' #Replacing punctations with space
                    break
                else:
                    x = ele
            filter_text2.append(x)
        filter_text3 = []
        text_B = []
        #Filtering all non-digital and non-alphabe elements in 
        #filter_text2 and store the filtered text in filter_text3

        for ele in filter_text2:
            it = re.finditer(r"\W",ele)
            text_A = []
            for match in it:
                a = match.group()
                text_A.append(a)

            if len(text_A)!=0:
                text_B = '' #Replacing with space
            else:
                text_B = ele #if it is alphabe, store to text_B
            filter_text3.append(text_B)

        filter_text4 = [word for word in filter_text3 if len(word)>3]
        clean_text = []
        #Remove space
        for s in filter_text4:
            if s!= '':
                clean_text.append(s)
        return clean_text
    
    def remove_similarwords(self,text):
        simple_text = []
        for word in text: 
            A = self.wnl.lemmatize(word, pos="v") #verbs
            simple_text.append(A)
        simple_text2 = []
        for word in simple_text:
            B = self.wnl.lemmatize(word, pos="n") #nouns
            simple_text2.append(B)
        simple_text3 = []
        for word in simple_text2:
            C = self.wnl.lemmatize(word, pos="a") #adj
            simple_text3.append(C)
        return simple_text3
    
    #Search keywords in each report and match it with default keywords lists
    def find_ppt_keywords (self,directory):
        text1 = self.extract_text_from_ppt(directory)
        text_filter1 = self.remove_stopwords(text1)
        text_simple1 = self.remove_similarwords(text_filter1)
        text_all = [text_simple1]
        dictionary = corpora.Dictionary(text_all) #build the dictionary
        corpus = [dictionary.doc2bow(text) for text in text_all]
        ldamodel = gensim.models.ldamodel.LdaModel(corpus, num_topics=1, id2word = dictionary, passes=10)
        text_result = ldamodel.print_topics(num_topics=1, num_words=20)
        ppt_keywords = self.remove_stopwords(text_result[0][1]) 

        return ppt_keywords
    
    def find_pdf_keywords (self,directory):
        text1 = self.extract_text_from_pdf(directory)
        text_filter1 = self.remove_stopwords(text1)
        text_simple1 = self.remove_similarwords(text_filter1)
        text_all = [text_simple1]
        dictionary = corpora.Dictionary(text_all)
        corpus = [dictionary.doc2bow(text) for text in text_all]
        ldamodel = gensim.models.ldamodel.LdaModel(corpus, num_topics=1, id2word = dictionary, passes=10)
        text_result = ldamodel.print_topics(num_topics=1, num_words=20)
        pdf_keywords = self.remove_stopwords(text_result[0][1])

        return pdf_keywords
    
    def score(self,file):
        ppt_compare = []
        result_len = []
        for key in keywords:
            for word1 in file:
                for word2 in key:
                    if word1==word2:
                        ppt_compare.append(word1)
            #In case that the length of keywords lists is different so score it by calculating the accounting ratio 
            result_len.append(len(ppt_compare)/len(key)) 
        
            #return 0 to run keyword1, keyword2, and keyword3, respectively.
            ppt_compare = []
        #print(filename)
        
        loc = result_len.index(max(result_len)) #Choose the longest list to identify the ppt/pdf's content
        if loc == 0:
            return "Finite Element Analysis"
        elif loc==1:
            return "Human Resource"
            #print("The topic of this report is Human Resource")
        elif loc==2:
            return "Economic"
            #print("The topic of this report is Economic")
            #if use print, the data type will be non-type so the main topic will not be displayed on the UI

    def extract_pptx_info(self,path_file):
        prs = Presentation(path_file)
        #Find Author
        author = prs.core_properties.author if prs.core_properties.author is not None else "SBL company"
        #Find modified date
        modified = prs.core_properties.modified.isoformat()
        return author,modified
    
    def extract_pdf_info(self,path_file):
        with open(path_file, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            info = reader.metadata
            author = info.author if info.author is not None else "SBL company"
            modified = info.modification_date
            return author, modified
    


      
    def run(self, file_path):
        result = []
        if file_path.endswith('.pptx'):
            PPT = self.find_ppt_keywords(file_path)
            author, modified= self.extract_pptx_info(file_path)
            result.append({'file_name': os.path.basename(file_path), 
                           'author': author, 
                           'date': modified,
                           'The main content of the report is': self.score(PPT)})

        elif file_path.endswith('.pdf'):
            PPT = self.find_pdf_keywords(file_path)
            author, subject = self.extract_pdf_info(file_path)
            result.append({'file_name': os.path.basename(file_path), 
                           'author': author, 
                           'date': subject, 
                           'The main content of the report is ': self.score(PPT)})

        return result
 

#Operation Test for FileAnalyzer without UI
'''
    def run(self, path_file):
        for filename in os.listdir(path_file):
            if filename.endswith('.pptx'):
                file_path = os.path.join(path_file, filename)            
                PPT = self.find_ppt_keywords (file_path)
                print(filename)
                #print(PPT)
                author,intro = self.extract_pptx_info(file_path)
                print(f'Author: {author}')
                print(f'Introduction: {intro}')
                self.score(PPT)           #compare keywords PPT with keywords
                print(type(self.score(PPT)))
                print("-------------------------")
                

            elif filename.endswith('.pdf'):
                file_path = os.path.join(file_path, filename)
                
                PPT = self.find_pdf_keywords (file_path)
                print(filename)            
                author, subject = self.extract_pdf_info(file_path)
                print(f'Author: {author}')
                print(f'Introduction: {subject}')
                self.score(PPT)           #compare keywords PDF with keywords
                print("-------------------------")
    

directory = os.path.dirname(os.path.abspath('FEA1.pptx'))
result = FileAnalyzer(directory)
result.run(directory)

FileAnalyzer(directory).run(directory)

#result.score(text)  
'''
     



